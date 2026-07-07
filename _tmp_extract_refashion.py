# -*- coding: utf-8 -*-
import sys, io
from pathlib import Path
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

base = Path(r'd:\users\Strophe\Documents\1-IA\La Clique Qui Recycle\JARVOS_recyclique\references\eco-organismes\partenaires\refashion')

import pypdf

for name in ['demande_conventionnement_DPAV_refashion_2025_vdef.pdf', 'Contrat_Type_DPAV_ESS_version_2024.pdf']:
    pdf = base / 'referentiels-officiels' / name
    print('=' * 80)
    print('PDF:', name)
    reader = pypdf.PdfReader(str(pdf))
    start = 14 if 'Contrat' in name else 0
    for i in range(start, len(reader.pages)):
        t = reader.pages[i].extract_text() or ''
        print(f'--- Page {i+1} ---')
        print(t)

try:
    import xlrd
    xls = base / 'referentiels-officiels' / "Matrice point d'apport.xls"
    wb = xlrd.open_workbook(str(xls))
    for sh in wb.sheets():
        print('=' * 80)
        print('XLS sheet:', sh.name, 'rows', sh.nrows, 'cols', sh.ncols)
        for r in range(min(50, sh.nrows)):
            row = [sh.cell_value(r, c) for c in range(sh.ncols)]
            print(row)
except Exception as e:
    print('XLS error:', e)

try:
    from pptx import Presentation
    pptx = base / 'divers' / '250203_WIKI_Prsentation_AMI_TLC_2025.pptx'
    prs = Presentation(str(pptx))
    print('=' * 80)
    print('PPTX slides:', len(prs.slides))
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            print(f'--- Slide {i+1} ---')
            print('\n'.join(texts))
except Exception as e:
    import traceback
    traceback.print_exc()
